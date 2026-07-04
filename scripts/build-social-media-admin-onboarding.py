from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "SOP" / "social-media-admin" / "Trivpass-Onboarding-Admin-Sosmed.pptx"

GREEN = RGBColor(0x1F, 0x4D, 0x3A)
TERRA = RGBColor(0xC2, 0x60, 0x3E)
OFF = RGBColor(0xFA, 0xF7, 0xF2)
INK = RGBColor(0x1D, 0x1D, 0x1B)
MUTED = RGBColor(0x6D, 0x6A, 0x62)
HAIR = RGBColor(0xD8, 0xD0, 0xC3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = 13.333, 7.5
M = 0.62


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             font="Inter", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             italic=False, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=15, color=INK, bullet=False,
                  font="Inter", spacing_after=5, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing_after)
        p.line_spacing = line_spacing
        if bullet:
            p.level = 0
            p.text = f"- {line}"
    return box


def add_rect(slide, x, y, w, h, fill=OFF, line=HAIR, radius=False):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    return shape


def add_rule(slide, y):
    line = slide.shapes.add_connector(1, Inches(M), Inches(y), Inches(W - M), Inches(y))
    line.line.color.rgb = HAIR
    line.line.width = Pt(1)


def add_footer(slide, no):
    add_text(slide, "Trivpass · Onboarding Admin Social Media", M, 7.08, 5.3, 0.2, 8.5, MUTED)
    add_text(slide, f"{no:02d}", 12.15, 7.08, 0.5, 0.2, 8.5, MUTED, align=PP_ALIGN.RIGHT)


def title(slide, kicker, headline, sub=None, no=None):
    add_text(slide, kicker.upper(), M, 0.46, 4.5, 0.3, 11, TERRA, bold=True)
    add_text(slide, headline, M, 0.86, 8.8, 0.55, 27, INK, bold=True, font="Fraunces")
    if sub:
        add_text(slide, sub, M, 1.45, 9.8, 0.42, 13.5, MUTED)
    add_rule(slide, 1.95)
    if no is not None:
        add_footer(slide, no)


def card(slide, x, y, w, h, head, body, num=None, fill=WHITE):
    add_rect(slide, x, y, w, h, fill=fill, line=HAIR)
    if num:
        add_text(slide, num, x + 0.22, y + 0.2, 0.35, 0.28, 13, TERRA, bold=True)
        tx = x + 0.62
        tw = w - 0.84
    else:
        tx = x + 0.24
        tw = w - 0.48
    add_text(slide, head, tx, y + 0.18, tw, 0.32, 14.5, INK, bold=True)
    add_multiline(slide, body if isinstance(body, list) else [body], tx, y + 0.64, tw, h - 0.86, 11.5, MUTED, spacing_after=3)


def add_table_like(slide, rows, x, y, w, h, col_ratio=(0.38, 0.62), header=False):
    row_h = h / len(rows)
    for i, row in enumerate(rows):
        yy = y + i * row_h
        fill = RGBColor(0xF3, 0xEE, 0xE5) if i == 0 and header else WHITE
        add_rect(slide, x, yy, w * col_ratio[0], row_h, fill=fill, line=HAIR)
        add_rect(slide, x + w * col_ratio[0], yy, w * col_ratio[1], row_h, fill=fill, line=HAIR)
        add_text(slide, row[0], x + 0.16, yy + 0.11, w * col_ratio[0] - 0.28, row_h - 0.16, 10.5, INK if i == 0 and header else GREEN, bold=True)
        add_text(slide, row[1], x + w * col_ratio[0] + 0.16, yy + 0.11, w * col_ratio[1] - 0.28, row_h - 0.16, 10.5, INK if i == 0 and header else MUTED)


def add_picture_fit(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        pic_w = w
        pic_h = w / img_ratio
        px = x
        py = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * img_ratio
        px = x + (w - pic_w) / 2
        py = y
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pic_w), Inches(pic_h))


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]


def new_slide(bg=OFF):
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


# 1
s = new_slide()
add_text(s, "trivpass", M, 0.52, 1.7, 0.3, 18, GREEN, bold=True, font="Fraunces")
add_text(s, ".", 1.92, 0.52, 0.2, 0.3, 18, TERRA, bold=True, font="Fraunces")
add_text(s, "ONBOARDING", M, 2.0, 4.2, 0.35, 20, TERRA, bold=True)
add_text(s, "Admin Social Media", M, 2.42, 7.6, 0.75, 40, INK, bold=True, font="Fraunces")
add_text(s, "Panduan kerja, brand, standar konten, dan dasar Meta Ads untuk Instagram & Facebook Trivpass.", M, 3.25, 7.3, 0.52, 15.5, MUTED)
add_text(s, "PT Taman Juada Legian · Operasional · v1.1 · Juni 2026", M, 6.72, 6, 0.24, 10.5, MUTED)

# 2
s = new_slide()
title(s, "Selamat datang di tim", "Peran kamu dalam satu kalimat", no=2)
add_text(s, "Kamu menjaga suara, tampilan, dan eksekusi dasar Trivpass di Instagram, Facebook, dan Meta Ads - supaya calon traveler percaya sebelum mereka memesan.", M, 2.25, 11.3, 0.62, 18, INK, bold=True)
card(s, M, 3.35, 3.75, 2.2, "Yang kamu pegang", ["Kalender konten", "Post, carousel, reel, story", "DM & komentar tahap awal", "Draft ads dan report bulanan"])
card(s, 4.8, 3.35, 3.75, 2.2, "Yang kamu jaga", ["Brand voice", "3 warna brand", "No-faces", "CTA dan klaim yang benar"])
card(s, 8.97, 3.35, 3.75, 2.2, "Cara kerja", ["Lapor ke Owner / Operasional", "Harga, klaim, budaya, dan ads wajib approve", "Konsistensi > kreasi bebas"])
add_footer(s, 2)

# 3
s = new_slide()
title(s, "Konteks brand", "Apa itu Trivpass", no=3)
add_text(s, "Travel agency untuk Bali yang menghapus calo antara traveler dan operator: driver yang kami vetting & kontrak sendiri, harga all-in, support Seminyak, dan pengalaman budaya yang dikurasi hati-hati. Pertahanannya ada di model operasi - bukan di software.", M, 2.25, 11.8, 0.75, 18, INK)
card(s, M, 3.55, 3.75, 1.65, "Real drivers", "Driver internal yang kami vetting & kontrak - bukan jemputan acak.", fill=WHITE)
card(s, 4.8, 3.55, 3.75, 1.65, "Real prices", "Harga all-in yang dirinci: driver, tiket, bensin, parkir, pajak.", fill=WHITE)
card(s, 8.97, 3.55, 3.75, 1.65, "No surprises", "Ops 24/7 + detail driver sebelum pickup. Traveler tahu siapa & jam berapa.", fill=WHITE)
add_footer(s, 3)

# 4
s = new_slide()
title(s, "Lensa strategi", "Mayoritas konten membangun trust", "3 kecemasan ini adalah panduan, bukan aturan kaku untuk setiap post.", no=4)
card(s, M, 2.35, 3.75, 2.3, "Takut kena scam", "Dijawab oleh driver roster: driver internal, di-vetting, bukan jemputan acak aplikasi gig.", "1")
card(s, 4.8, 2.35, 3.75, 2.3, "Takut kemahalan", "Dijawab oleh harga all-in yang dirinci: driver, tiket, bensin, parkir, pajak, biaya layanan.", "2")
card(s, 8.97, 2.35, 3.75, 2.3, "Takut booking gagal", "Dijawab oleh ops 24/7 + email perkenalan driver: siapa yang menjemput, dan jam berapa.", "3")
add_rect(s, M, 5.15, 12.1, 0.72, fill=RGBColor(0xF3, 0xEE, 0xE5), line=HAIR)
add_text(s, "Gunakan lensa ini terutama untuk trust-building, product post, FAQ, dan ads. Konten budaya, route idea, atau seasonal education boleh lebih natural selama tetap spesifik dan tidak overclaim.", M + 0.25, 5.33, 11.6, 0.32, 12.5, GREEN, bold=True)
add_footer(s, 4)

# 5
s = new_slide()
title(s, "Cara Trivpass bicara", "Brand voice", no=5)
add_text(s, "Seperti teman lokal yang paham teknologi, sudah 10 tahun di bidang ini, dan tidak akan menjual berlebihan - tenang, spesifik, transparan, sesekali datar.", M, 2.18, 11.4, 0.55, 18, INK, bold=True)
card(s, M, 3.1, 5.75, 1.15, "1 · Mulai dari kebutuhan traveler", "Mulai dari pertanyaan, situasi, atau risiko traveler. Cerita Trivpass jadi pendukung.")
card(s, 6.95, 3.1, 5.75, 1.15, "2 · Spesifik - angka & nama", "Pakai jam, durasi, harga, wilayah, dan proses. Kata sifat harus kalah dari detail.")
card(s, M, 4.55, 5.75, 1.15, "3 · Itemize", "Harga, inclusions, exclusions, dan refund window dirinci. Transparansi itu layout.")
card(s, 6.95, 4.55, 5.75, 1.15, "4 · Jangan jual yang tak bisa dikirim", "Bukan 'magical sunrise', tapi 'pickup 03:30, summit 06:15, turun 09:00'.")
add_footer(s, 5)

# 6
s = new_slide()
title(s, "Disiplin bahasa", "Kata yang dipakai - dan yang dilarang", no=6)
add_table_like(s, [
    ("Pakai ini", "Jangan pernah"),
    ("Verified driver", '"Book now" -> pakai CTA Trivpass'),
    ("All-in price", '"Premium" / "Luxury" / "World-class"'),
    ("From IDR ...", '"Authentic" / "Hidden gem" / "Bucket list"'),
    ("IDR 100.000 ke host community", '"Once-in-a-lifetime" / "Don\'t miss out"'),
    ("24/7 ops in Seminyak", '"#1 in Bali" / "Hassle-free" / klaim tanpa bukti'),
], M, 2.25, 12.1, 3.9, header=True)
add_footer(s, 6)

# 7
s = new_slide()
title(s, "Sistem visual", "Tiga warna. Dua font. Tanpa kompromi.", no=7)
for x, color, name, use in [
    (M, GREEN, "#1F4D3A\nJungle Green", "Logo, CTA, status verified"),
    (4.8, TERRA, "#C2603E\nTerracotta", "Harga, rating, scarcity - sinyal, bukan tombol"),
    (8.97, OFF, "#FAF7F2\nWarm Off-White", "Background dasar semua halaman"),
]:
    add_rect(s, x, 2.35, 3.75, 1.15, fill=color, line=HAIR)
    add_text(s, name, x + 0.22, 2.58, 3.25, 0.44, 18, WHITE if color != OFF else INK, bold=True)
    add_text(s, use, x + 0.22, 3.75, 3.35, 0.36, 12, MUTED)
add_text(s, "Kuota: ~85% canvas + teks · <=3% Jungle Green · <=2% Terracotta. Tanpa warna keempat, tanpa hitam/putih murni, tanpa gradient/bayangan dekoratif.", M, 4.72, 11.6, 0.42, 13.5, INK)
add_text(s, "Fraunces untuk judul/headline · Inter untuk body, harga, tanggal, semua angka.", M, 5.55, 10.8, 0.36, 16, GREEN, bold=True)
add_footer(s, 7)

# 8
s = new_slide()
title(s, "Aturan terkunci", "Lima hal yang tidak bisa ditawar", no=8)
rules = [
    ("Logo di paling atas", "Mark jungle-green + 'trivpass' serif kecil + titik terracotta."),
    ("Tanpa wajah", "Orang dari belakang / siluet / tangan boleh. Tidak ada pose influencer."),
    ("Hanya 3 warna brand", "Off-white, jungle-green, terracotta. Tidak ada warna keempat."),
    ("Foto hangat, tidak glossy", "Sedikit desaturasi. Tanpa HDR, gradient, atau shadow dekoratif."),
    ("AI image sementara boleh", "Harus terlihat nyata, dokumenter, no-faces. Target jangka panjang: foto asli."),
]
for i, (h, b) in enumerate(rules):
    y = 2.22 + i * 0.86
    add_text(s, str(i + 1), M, y, 0.35, 0.28, 16, TERRA, bold=True)
    add_text(s, h, 1.05, y, 4.0, 0.28, 14, INK, bold=True)
    add_text(s, b, 5.05, y, 7.3, 0.36, 12, MUTED)
add_footer(s, 8)

# 9
s = new_slide()
title(s, "Apa yang kita posting", "Enam pilar konten", no=9)
pillars = [
    ("Bali Travel Education", "Tips praktis, checklist, hal yang harus dihindari saat merencanakan Bali."),
    ("Trust Building", "Driver vetting, ops 24/7, refund - bukti, bukan klaim."),
    ("Brand Introduction", "Siapa Trivpass, model all-in, kenapa beda dari OTA."),
    ("Tour Packages", "Spotlight tour: rute, jam, yang termasuk, harga hero."),
    ("Activities & Attractions", "Tiket atraksi, harga transparan, add-on driver opsional."),
    ("Cultural Literacy", "Penjelasan upacara & istilah Bali - edukasi, bukan dorongan booking."),
]
for i, (h, b) in enumerate(pillars):
    x = M + (i % 3) * 4.18
    y = 2.28 + (i // 3) * 1.72
    card(s, x, y, 3.75, 1.35, h, b)
add_footer(s, 9)

# 10
s = new_slide()
title(s, "Disiplin ajakan", "CTA: lembut saat awareness, langsung saat produk", no=10)
add_rect(s, M, 2.3, 12.1, 0.72, fill=GREEN, line=GREEN)
add_text(s, "CTA utama produk: \"Secure Your Driver & Tickets\"", M + 0.25, 2.49, 11.5, 0.28, 17, WHITE, bold=True)
card(s, M, 3.55, 5.75, 1.75, "Soft CTA - awareness", ["Save post ini", "Follow untuk tips Bali", "DM pertanyaan kamu", "Comment bulan rencana liburanmu"])
card(s, 6.95, 3.55, 5.75, 1.75, "Direct CTA - produk", ["Atraksi/tiket -> Secure Your Tickets", "Tour + driver -> Secure Your Driver & Tickets", "Event budaya -> Reserve Your Spot bila live"])
add_text(s, "Tidak pernah: \"Book now\" / \"Reserve\" / \"Get started\" di copy Trivpass, kecuali label platform Meta memaksa.", M, 5.85, 11.6, 0.32, 12.5, TERRA, bold=True)
add_footer(s, 10)

# 11
s = new_slide()
title(s, "Alur produksi", "Tiga langkah kerja", no=11)
card(s, M, 2.35, 3.75, 2.25, "Susun perencanaan", "Kembangkan content plan bulan berjalan jadi rencana: tanggal, format, pilar, hook, CTA, dan status.", "1")
card(s, 4.8, 2.35, 3.75, 2.25, "Preview & approval", "Kirim preview ke Owner. Harga, klaim, budaya, produk baru, dan ads selalu tunggu approval.", "2")
card(s, 8.97, 2.35, 3.75, 2.25, "Revisi & posting", "Satu round revisi, lalu posting/jadwalkan di IG utama + cross-post Facebook.", "3")
add_text(s, "File post disimpan di content/{tahun}/{bulan}/{tgl}-{judul}/. Status: draft -> needs asset -> ready -> blocked -> posted.", M, 5.28, 11.8, 0.35, 12.5, GREEN, bold=True)
add_footer(s, 11)

# 12
s = new_slide()
title(s, "Meta Ads", "Ruang lingkup admin", "Admin boleh menjalankan eksekusi dasar, tapi keputusan besar tetap milik Owner.", no=12)
card(s, M, 2.35, 3.75, 2.3, "Boleh dilakukan", ["Draft copy & creative ads", "Audience idea", "Campaign plan", "Screenshot hasil & report bulanan"])
card(s, 4.8, 2.35, 3.75, 2.3, "Butuh approval", ["Campaign sebelum live", "Budget & perubahan budget", "Target market", "Harga, promo, klaim, landing page"])
card(s, 8.97, 2.35, 3.75, 2.3, "Tidak boleh", ["Angle diskon murahan", "Countdown palsu", "Overclaim cultural events", "Ubah campaign besar tanpa izin"])
add_text(s, "Produk aman untuk ads awal: Tour Packages dan Activities & Attractions. Cultural Events hanya dipromosikan sebagai edukasi sampai status live dikonfirmasi.", M, 5.25, 11.9, 0.46, 12.5, TERRA, bold=True)
add_footer(s, 12)

# 13
s = new_slide()
title(s, "Ads yang dipelajari", "Roadmap iklan Trivpass", no=13)
rows = [
    ("Fase", "Yang dipelajari admin"),
    ("0 · Foundation", "Business Suite, pixel/dataset, domain, event, payment, UTM. Tidak spend sebelum setup jelas."),
    ("1 · Warm Launch", "Engagement/Traffic. Creative: brand intro, driver vetting, transparent pricing. Cari angle terbaik."),
    ("2 · Trip Consideration", "Traffic/Leads. Creative: Ubud, Mt. Batur, Uluwatu, Activities & Attractions. Ukur inquiry."),
    ("3 · Conversion-ready test", "Retargeting/Sales bila tracking siap. Pakai all-in pricing dan inclusions, bukan promo palsu."),
]
add_table_like(s, rows, M, 2.28, 12.1, 3.8, col_ratio=(0.28, 0.72), header=True)
add_text(s, "Campaign naming: TP_[Objective]_[Market]_[Audience]_[CreativeAngle]_[YYYYMM]", M, 6.35, 11.5, 0.26, 11.5, GREEN, bold=True)
add_footer(s, 13)

# 14
s = new_slide()
title(s, "Target skill ads", "Format iklan yang harus dipelajari", "Tujuannya: naik dari boost content ke Ads Manager dan product discovery ads.", no=14)
add_text(s, "Benchmark visual: product card, carousel, story, grid. Gaya Trivpass tetap calm: all-in price, verified driver, tanpa urgency palsu.", M, 1.73, 11.6, 0.18, 8.5, TERRA, bold=True)
rows = [
    ("Format", "Target kemampuan admin"),
    ("TPD-1 Tour Carousel", "Manual Carousel Ads. Tiap card bisa membawa produk/link berbeda: Mt. Batur, Uluwatu, Waterbom, ATV."),
    ("TPD-2 Single Product Story", "Story/Reels vertical 9:16 untuk 1 produk. Foto/video kuat, harga all-in, CTA Trivpass."),
    ("TPD-3 Activities Grid", "Collection-style grid untuk beberapa produk. Cocok untuk Activities & Attractions dan family trips."),
    ("TPD-4 Website Traffic", "Ads direct ke website: halaman produk, kategori, atau landing page campaign. Fokus landing page views."),
    ("TPD-5 Retargeting Reminder", "Ads untuk IG engagers, website visitors, atau cart/checkout abandoners bila tracking siap."),
    ("Catalog readiness", "Belajar struktur catalog/product feed: nama, harga, gambar, URL, kategori, availability."),
]
add_table_like(s, rows, M, 2.18, 12.1, 4.45, col_ratio=(0.34, 0.66), header=True)
add_footer(s, 14)

# 15
s = new_slide()
title(s, "Di mana format ini ada", "Mapping ke Meta Ads Manager", "Nama menu Meta bisa berubah, tapi konsepnya tetap: objective -> placement -> format -> destination.", no=15)
rows = [
    ("Format", "Di Meta Ads Manager"),
    ("TPD-1 Tour Carousel", "Campaign: Traffic/Leads/Sales. Ad level: Format = Carousel. Isi card manual; tiap card bisa punya URL berbeda."),
    ("TPD-2 Single Product Story", "Ad set: pilih placement Instagram Stories/Reels. Ad level: Single image/video atau carousel vertical 9:16."),
    ("TPD-3 Activities Grid", "Ad level: Collection, atau Sales + Catalog/Product Set. Cocok untuk membuka katalog/Instant Experience."),
    ("TPD-4 Website Traffic", "Campaign: Traffic atau Sales. Conversion location: Website. Destination: product page, category page, atau checkout-ready landing page."),
    ("TPD-5 Retargeting", "Ad set: Custom Audience. Contoh: IG engagers, website visitors, cart/checkout abandoners."),
    ("Catalog readiness", "Commerce Manager / Business assets: buat Catalog, product feed, product set. Dipakai untuk Advantage+ catalog ads."),
]
add_table_like(s, rows, M, 2.18, 12.1, 4.45, col_ratio=(0.32, 0.68), header=True)
add_footer(s, 15)

# 16
s = new_slide()
title(s, "Visual benchmark", "Referensi GetYourGuide", "Pakai ini untuk membaca struktur ads, bukan untuk meniru gaya copy atau urgency-nya.", no=16)
img1 = ROOT / "SOP/social-media-admin/reference-ads/getyourguide-carousel-feed.png"
img2 = ROOT / "SOP/social-media-admin/reference-ads/getyourguide-story-product.png"
img3 = ROOT / "SOP/social-media-admin/reference-ads/getyourguide-story-grid.png"
for y, path, eyebrow, label, body in [
    (2.22, img1, "CONTOH 1", "Feed carousel / product card", "Banyak kartu produk. Setiap card bisa membawa link berbeda ke halaman produk."),
    (3.58, img2, "CONTOH 2", "Story single product", "Satu produk ditonjolkan dalam format 9:16. Cocok untuk tour atau attraction hero."),
    (4.94, img3, "CONTOH 3", "Story grid / collection", "Beberapa produk ditampilkan sebagai grid. Cocok untuk Activities & Attractions."),
]:
    add_rect(s, M, y, 12.1, 1.04, fill=WHITE, line=HAIR)
    add_picture_fit(s, path, M + 0.16, y + 0.12, 1.18, 0.8)
    add_text(s, eyebrow, M + 1.58, y + 0.15, 1.1, 0.18, 8.2, TERRA, bold=True)
    add_text(s, label, M + 1.58, y + 0.39, 3.2, 0.22, 12.2, GREEN, bold=True)
    add_text(s, body, M + 5.0, y + 0.3, 6.7, 0.34, 11.1, MUTED)
add_rect(s, M, 6.34, 12.1, 0.36, fill=RGBColor(0xF3, 0xEE, 0xE5), line=HAIR)
add_text(s, "Adaptasi Trivpass: all-in price, verified/optional driver, rute jelas, CTA Trivpass, tanpa urgency palsu.", M + 0.25, 6.45, 11.55, 0.15, 9.8, TERRA, bold=True)
add_footer(s, 16)

# 17
s = new_slide()
title(s, "Target konten bulanan", "Output minimum yang harus dipenuhi", no=17)
rows = [
    ("Kategori", "Target & tujuan"),
    ("Feed/Reel utama", "12 konten per bulan. Campuran edukasi, trust, produk, FAQ."),
    ("Story set", "8-12 set per bulan. Poll, FAQ, behind-the-scenes, repost, trip prompt."),
    ("Bali Travel Education", "4 konten. Target: saves dan shares."),
    ("Trust Building", "3 konten. Target: komentar/DM yang menunjukkan rasa percaya."),
    ("Tour Packages", "2 konten. Target: itinerary views atau inquiry."),
    ("Activities & Attractions", "2 konten. Target: clicks, saves, atau DM."),
    ("Brand/FAQ/BTS", "1-2 konten. Target: membuat Trivpass mudah dipahami."),
]
add_table_like(s, rows, M, 2.18, 12.1, 4.45, col_ratio=(0.34, 0.66), header=True)
add_footer(s, 17)

# 18
s = new_slide()
title(s, "Indikator bulanan", "Konten dan ads dinilai per bulan", no=18)
card(s, M, 2.25, 5.75, 2.8, "Organic content", ["Jumlah konten tayang vs rencana", "Top 3 konten: saves, shares, reach, engagement", "DM/comment inquiry", "Pertanyaan traveler yang berulang", "Konten lemah + pelajaran bulan depan"])
card(s, 6.95, 2.25, 5.75, 2.8, "Meta Ads", ["Total spend", "Reach, clicks, landing page views", "CPC / cost per landing page view", "Website inquiry / checkout signal", "Cost per inquiry atau checkout start", "Creative, audience, dan format terbaik"])
add_rect(s, M, 5.55, 12.1, 0.68, fill=RGBColor(0xF3, 0xEE, 0xE5), line=HAIR)
add_text(s, "Jangan mengejar viral kosong. Untuk Trivpass, indikator sehat adalah trust signal + inquiry signal.", M + 0.25, 5.74, 11.55, 0.28, 12.5, GREEN, bold=True)
add_footer(s, 18)

# 19
s = new_slide()
title(s, "Sebelum tayang", "Yang wajib approve Owner dulu", no=19)
card(s, M, 2.25, 5.75, 1.35, "Harga & angka", "Semua harga, diskon, rating, budget ads, dan spend. Jangan pernah mengarang figure.")
card(s, 6.95, 2.25, 5.75, 1.35, "Klaim & janji", "Driver, refund, partner, ketersediaan, tracking, dan campaign promise harus bisa dibuktikan.")
card(s, M, 3.95, 5.75, 1.35, "Konten sensitif budaya", "Upacara, istilah Bali, foto tempat sakral. Hormat sebagai tamu, bukan penonton.")
card(s, 6.95, 3.95, 5.75, 1.35, "Apa pun yang baru", "Produk baru, identitas driver, partner, tanggal event, landing page baru, campaign baru.")
add_text(s, "Aturan emas: jangan pernah mengarang foto, harga, driver, partner, ketersediaan event, tanggal, atau hasil ads. Tandai yang belum pasti - jangan tebak.", M, 5.86, 11.8, 0.42, 12.5, TERRA, bold=True)
add_footer(s, 19)

# 20
s = new_slide()
title(s, "Cheat sheet", "Kalau ragu, ingat enam ini", no=20)
cheats = [
    ("Trust dulu", "Mayoritas konten membangun percaya: anti-scam, harga jelas, booking jalan."),
    ("Spesifik > kata sifat", "Angka, nama, jam. Hapus kata sifat yang tak bisa jadi detail."),
    ("3 warna, no faces", "Off-white, jungle-green, terracotta. Tanpa wajah, tanpa gloss."),
    ("CTA yang benar", "Secure Your Driver & Tickets. Tidak pernah Book now di copy."),
    ("Ads perlu approval", "Budget, audience, klaim, creative, dan landing page tidak jalan tanpa izin."),
    ("Jangan mengarang", "Harga, driver, partner, tanggal, event, hasil ads - approve dulu ke Owner."),
]
for i, (h, b) in enumerate(cheats):
    x = M + (i % 3) * 4.18
    y = 2.2 + (i // 3) * 1.75
    card(s, x, y, 3.75, 1.38, h, b, str(i + 1))
add_text(s, "Real drivers. Real prices. No surprises.", M, 6.25, 11.5, 0.35, 18, GREEN, bold=True, font="Fraunces")
add_footer(s, 20)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
